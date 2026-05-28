"""Genera la presentación PowerPoint del QA de la integración Connecteam -> Odoo.

Uso:  /Users/dacm/we/.venv/bin/python qa/build_presentation.py
Salida: qa/Presentacion_QA_Integracion.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paleta
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x8B, 0x9E)
LIGHT = RGBColor(0xEE, 0xF1, 0xF5)
ROW = RGBColor(0xF7, 0xF9, 0xFB)
DARK = RGBColor(0x24, 0x2B, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x57)
RED = RGBColor(0xB3, 0x3A, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _no_border(shape):
    shape.line.fill.background()


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def slide():
    return prs.slides.add_slide(BLANK)


def title_bar(s, title, kicker=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    _fill(bar, NAVY)
    _no_border(bar)
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if kicker:
        p2 = tf.add_paragraph()
        p2.text = kicker
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0xBF, 0xD2, 0xDE)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SW, Inches(0.06))
    _fill(strip, ACCENT)
    _no_border(strip)


def bullets(s, items, left=Inches(0.7), top=Inches(1.5), width=None, height=None, size=16):
    width = width or (SW - Inches(1.4))
    height = height or (SH - top - Inches(0.4))
    box = s.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        marker = "" if level == 0 else "– "
        p.text = ("" if level == 0 else marker) + text
        p.font.size = Pt(size - 2 * level)
        p.font.color.rgb = DARK if level == 0 else RGBColor(0x44, 0x4C, 0x55)
        p.font.bold = (level == 0)
        p.space_after = Pt(7)
        p.space_before = Pt(2)
        if level == 0:
            # vinaza de color
            r = p.runs[0]
            r.font.color.rgb = NAVY
    return box


def table(s, data, top=Inches(1.55), left=Inches(0.6), col_widths=None,
          font=12, header=True):
    rows, cols = len(data), len(data[0])
    width = SW - Inches(1.2)
    height = Inches(0.4) * rows
    gtab = s.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtab.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = gtab.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font)
            if header and r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                para.font.color.rgb = WHITE
                para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW if r % 2 else LIGHT
                para.font.color.rgb = DARK
    return gtab


def section(title, n):
    s = slide()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _fill(bg, NAVY); _no_border(bg)
    num = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(3), Inches(1))
    p = num.text_frame.paragraphs[0]
    p.text = f"{n:02d}"
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = ACCENT
    tb = s.shapes.add_textbox(Inches(0.85), Inches(3.5), SW - Inches(1.6), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = WHITE
    tb.text_frame.word_wrap = True


def footer(s, txt="QA Integración Connecteam → Odoo"):
    box = s.shapes.add_textbox(Inches(0.5), SH - Inches(0.4), SW - Inches(1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = txt
    p.font.size = Pt(9); p.font.color.rgb = RGBColor(0x9A, 0xA4, 0xAE)


# ---------------------------------------------------------------- 1. PORTADA
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH); _fill(bg, NAVY); _no_border(bg)
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.05), SW, Inches(0.08))
_fill(strip, ACCENT); _no_border(strip)
t = s.shapes.add_textbox(Inches(0.9), Inches(2.2), SW - Inches(1.8), Inches(1.8))
p = t.text_frame.paragraphs[0]
p.text = "Aseguramiento de Calidad (QA)"
p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
t.text_frame.word_wrap = True
sub = s.shapes.add_textbox(Inches(0.9), Inches(4.3), SW - Inches(1.8), Inches(1.5))
p = sub.text_frame.paragraphs[0]
p.text = "Integración Connecteam → Odoo"
p.font.size = Pt(22); p.font.color.rgb = ACCENT
p2 = sub.text_frame.add_paragraph()
p2.text = "Pruebas automatizadas, defectos encontrados y gestión manual de excepciones"
p2.font.size = Pt(15); p2.font.color.rgb = RGBColor(0xC7, 0xD3, 0xDD)
p3 = sub.text_frame.add_paragraph()
p3.text = "Mayo 2026"
p3.font.size = Pt(12); p3.font.color.rgb = RGBColor(0x8C, 0x9A, 0xA6)

# ---------------------------------------------------------------- 2. AGENDA
s = slide(); title_bar(s, "Agenda")
bullets(s, [
    "Contexto: qué hace la integración y por qué es delicada",
    "El problema que resolvimos con QA",
    "Estrategia: pirámide de pruebas y el OdooSpy",
    "Resultados y cobertura por módulo",
    "Validación contra el Odoo real (incluye movimiento de equipos)",
    "Defectos encontrados (y corregidos)",
    "Gestión manual de excepciones: el inbox por etiqueta y por tipo de trabajo",
    "Puntos críticos del flujo de datos y brechas pendientes",
], size=17)
footer(s)

# ---------------------------------------------------------------- 3. CONTEXTO
s = slide(); title_bar(s, "Contexto: el \"secretario virtual\"", "Qué hace la integración")
bullets(s, [
    "Pipeline programado que toma los formularios que llenan los técnicos en terreno (Connecteam)",
    "Crea y actualiza órdenes de mantención en Odoo, mueve equipos de ubicación, genera informes PDF y deja avisos",
    "Cinco tipos de trabajo:",
    ("MC = Correctiva · CF = Configuración · MP = Preventiva", 1),
    ("I = Instalación · R = Reemplazo/Extracción", 1),
    "El núcleo es processor.py: ~4.500 líneas en una sola función",
    "Si el secretario se equivoca, hoy no avisa: sigue en silencio",
], size=16)
footer(s)

# ---------------------------------------------------------------- 4. PROBLEMA
s = slide(); title_bar(s, "El problema", "Por qué hacía falta QA")
bullets(s, [
    "Sin pruebas y sin linter: cada cambio era a ciegas",
    "Escribe en Odoo vía XML-RPC: los errores tienen efectos reales",
    "Falla en silencio: el patrón try/except + continue traga el error y continúa",
    ("Una OT puede \"procesarse\" y no crear nada, sin que nadie se entere", 1),
    "Estado de deduplicación e IDs hardcodeados que son carga útil, no configuración",
    "4.500 líneas de lógica de negocio acoplada y difícil de auditar",
], size=16)
footer(s)

# ---------------------------------------------------------------- 5. SECCIÓN
section("Estrategia de QA", 1)

# ---------------------------------------------------------------- 6. PIRÁMIDE
s = slide(); title_bar(s, "Pirámide de pruebas", "De barato y rápido a caro y realista")
bullets(s, [
    "L0 Estático: imports válidos, inventario de IDs hardcodeados",
    "L1 Unitario: funciones puras (parsing de respuestas, deduplicación)",
    "L2 Componente: el pipeline completo con un doble de prueba (OdooSpy), sin red",
    ("Aquí está el grueso de la cobertura: todas las ramas de decisión", 1),
    "L3 Integración: contra el Odoo de pruebas real (staging)",
    "Regla de oro: nunca aceptar \"no dio error\". Siempre afirmar el efecto concreto",
], size=16)
footer(s)

# ---------------------------------------------------------------- 7. SPY
s = slide(); title_bar(s, "El OdooSpy", "Cómo probamos sin tocar Odoo ni esperar la red")
bullets(s, [
    "Un doble de prueba que imita la interfaz de OdooClient",
    "Registra cada llamada (create / write / search / message_post...)",
    "Devuelve respuestas programables: le decimos \"el equipo existe y está aquí\"",
    "Permite simular cualquier estado de Odoo y luego revisar qué órdenes dio el secretario",
    "Resultado: se cubren todas las decisiones posibles, de forma rápida y determinista",
], size=16)
footer(s)

# ---------------------------------------------------------------- 8. SECCIÓN
section("Resultados y cobertura", 2)

# ---------------------------------------------------------------- 9. RESULTADOS
s = slide(); title_bar(s, "Resultados: 77 pruebas verdes", "Estado verificado · evidencia en RESULTADOS.md")
table(s, [
    ["Nivel", "Pruebas", "Qué cubre"],
    ["L1 Unitario", "12", "Parsing de respuestas + deduplicación (DB aislada)"],
    ["L2 Componente", "46", "Todas las ramas de los 5 módulos (con OdooSpy)"],
    ["L3 Integración", "19", "4 de solo lectura + 15 E2E que escriben en el staging"],
    ["TOTAL", "77", "Todas en verde"],
], top=Inches(2.0), col_widths=[3.0, 1.6, 7.5], font=14)
footer(s)

# ---------------------------------------------------------------- 10. COBERTURA MÓDULOS
s = slide(); title_bar(s, "Cobertura por módulo (L2)", "Cada tipo de trabajo, todas sus ramas")
table(s, [
    ["Módulo", "Tests", "Lógica distintiva probada"],
    ["MC  Correctiva", "12", "Interruptor: vincular a activa vs crear"],
    ["CF  Configuración", "9", "Proximidad temporal + archivar anteriores"],
    ["MP  Preventiva", "9", "Proximidad + archivado + 'sin plan'"],
    ["I   Instalación", "9", "Primera activa + escribe ubicación del equipo"],
    ["R   Reemplazo", "7", "Bifásico (E/I) + mueve equipo + Metrocal"],
], top=Inches(1.9), col_widths=[3.2, 1.4, 7.5], font=13)
bullets(s, [
    "Ramas comunes cubiertas: S/N no encontrado, punto inexistente, validación de ubicación, crear vs actualizar, operativo Sí/No",
], top=Inches(5.4), size=13)
footer(s)

# ---------------------------------------------------------------- 11. INTEGRACIÓN
s = slide(); title_bar(s, "Validación contra el Odoo real", "Staging: escribe de verdad")
bullets(s, [
    "19 pruebas contra el test-Odoo (instancia de staging)",
    "4 de solo lectura: autenticación + verificación de IDs hardcodeados (partners, Metrocal 2, 593/594)",
    "15 E2E de escritura que ejecutan el pipeline de punta a punta:",
    ("Camino feliz por módulo + movimiento real de ubicación (I → punto; R → 593 / 594 / punto)", 1),
    ("Enrutamiento de excepciones al inbox (S/N no encontrado, Punto no existe)", 1),
    ("operativo=No (stage 3 + adjunto), vincular a existente, proximidad + archivado", 1),
    "El reporte vincula los 42 registros reales creados con la prueba que los originó",
], size=14)
footer(s)

# ---------------------------------------------------------------- 12. DEFECTOS
s = slide(); title_bar(s, "Defectos encontrados", "El valor real del QA")
b = bullets(s, [
    "OBS-10 (CORREGIDO): NameError en Instalación-no-operativa",
    ("El request se creaba pero se perdían en silencio el registro de éxito y el PDF", 1),
    ("Eran dos instancias del mismo bug; ambas corregidas", 1),
    "Otras 9 observaciones documentadas con su caso testigo:",
    ("Punto de dos dígitos mal detectado, posible doble conversión de zona horaria", 1),
    ("Columnas duplicadas en R, mapas de IDs que difieren prod/test, 'sin plan' de MP", 1),
], size=15)
footer(s)

# ---------------------------------------------------------------- 13. SECCIÓN
section("Gestión manual de excepciones", 3)

# ---------------------------------------------------------------- 14. INBOX ORIGENES
s = slide(); title_bar(s, "El inbox y los tres orígenes", "Cuánto debe hacer el operario")
bullets(s, [
    "Cuando el secretario no puede actuar, deja un registro en x_inbox_integracion",
    "Cada registro tiene una etiqueta (qué pasó) y un origen (cuánto hay que intervenir):",
    ("A · Automática: el secretario resolvió todo. Sin acción", 1),
    ("M · Manual: el secretario NO pudo. El operario ejecuta lo que faltó", 1),
    ("N · Notificación: el secretario ya actuó. El operario solo valida", 1),
    "Followers notificados automáticamente; 'Creación en espera' avisa a Juan",
], size=16)
footer(s)

# ---------------------------------------------------------------- 15. FLUJOS POR ETIQUETA
s = slide(); title_bar(s, "Flujos manuales por etiqueta")
table(s, [
    ["Etiqueta", "Origen", "Acción del operario"],
    ["S/N no encontrado", "M", "Crear el equipo (o vincular) y el evento"],
    ["Creación en espera", "M", "Esperar transferencia; Juan crea el equipo"],
    ["Punto no existe en sistema", "M", "Solicitar/crear el punto o vincular al correcto"],
    ["Cambio de ubicación", "N", "Validar punto y fecha en el equipo"],
    ["Sin evento de instalación", "N", "Validar por qué el equipo no tenía instalación"],
], top=Inches(1.7), col_widths=[3.8, 1.3, 7.0], font=13)
footer(s)

# ---------------------------------------------------------------- 16. POR TIPO DE TRABAJO
s = slide(); title_bar(s, "Acciones manuales por tipo de trabajo", "La distinción clave")
table(s, [
    ["Tipo", "Solicitud (vincular vs crear)", "¿Mueve el equipo?"],
    ["MC", "Interruptor (vincular activa o crear)", "No"],
    ["CF / MP", "Proximidad temporal + archivar viejas", "No"],
    ["I", "Primera activa o crear", "Sí → al punto"],
    ["R", "Extracción / Instalación / Calibración", "Sí → 593 / 594 / punto"],
], top=Inches(1.7), col_widths=[1.6, 6.6, 3.9], font=13)
bullets(s, [
    "MC / CF / MP terminan en la solicitud y nunca tocan el equipo",
    "I y R tienen un paso obligatorio extra sobre maintenance.equipment (la ubicación). En R depende del subtrabajo y la calibración suma a Metrocal",
], top=Inches(4.7), size=14)
footer(s)

# ---------------------------------------------------------------- 17. SECCIÓN
section("Puntos críticos del flujo de datos", 4)

# ---------------------------------------------------------------- 18. PUNTOS CRÍTICOS
s = slide(); title_bar(s, "Puntos críticos del flujo de datos", "Dónde se rompe la cadena")
bullets(s, [
    "Convención de columnas del formulario: contrato implícito, sin validación de esquema",
    "Llaves de cruce: número de serie y nombre del punto deben coincidir exacto",
    "Errores silenciosos: hoy un fallo no detiene ni alerta",
    "Deduplicación e IDs hardcodeados (difieren prod/test)",
    "Selección de solicitud y movimiento de ubicación (estado irreversible en I/R)",
    "Triaje del inbox: el fallback humano debe ejecutarse y replicar el estado exacto",
], size=15)
footer(s)

# ---------------------------------------------------------------- 19. BRECHAS
s = slide(); title_bar(s, "Brechas y recomendaciones", "Para garantizar el flujo en producción")
table(s, [
    ["Brecha", "Recomendación"],
    ["El formulario puede cambiar y romper el parsing en silencio", "Validar el esquema del formulario antes de procesar"],
    ["El sistema falla en silencio", "Monitoreo / alerta cuando una OT no produce efecto"],
    ["IDs hardcodeados difieren prod/test", "Revalidar IDs en cada promoción o migración"],
    ["El inbox depende de revisión humana", "Triaje disciplinado siguiendo el runbook"],
], top=Inches(1.9), col_widths=[5.6, 6.5], font=13)
bullets(s, [
    "La lógica de negocio ya está blindada por las 68 pruebas. Las dos brechas prioritarias son: validar el formulario y tener observabilidad.",
], top=Inches(5.3), size=13)
footer(s)

# ---------------------------------------------------------------- 20. CIERRE
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH); _fill(bg, NAVY); _no_border(bg)
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SW, Inches(0.08))
_fill(strip, ACCENT); _no_border(strip)
t = s.shapes.add_textbox(Inches(0.9), Inches(2.3), SW - Inches(1.8), Inches(3.5))
tf = t.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "En síntesis"
p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
for txt in [
    "77 pruebas automatizadas en 3 niveles, todas en verde",
    "Un bug oculto encontrado y corregido",
    "Validado contra el Odoo real, incluido el movimiento de equipos",
    "Manejo manual de excepciones documentado por etiqueta y por tipo de trabajo",
    "Dos prioridades para producción: validar el formulario y observabilidad",
]:
    pp = tf.add_paragraph(); pp.text = "•  " + txt
    pp.font.size = Pt(16); pp.font.color.rgb = RGBColor(0xD7, 0xE0, 0xE8)
    pp.space_after = Pt(8)

OUT = Path(__file__).resolve().parent / "Presentacion_QA_Integracion.pptx"
prs.save(str(OUT))
print(f"OK -> {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
